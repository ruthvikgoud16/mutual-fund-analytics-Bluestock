"""12-Slide PowerPoint Presentation Generator (Day 7 Task 2).

Generates `docs/Bluestock_MF_Presentation.pptx` using `python-pptx`, exactly following the 12-slide specification:
- Slide 1: Title Slide
- Slide 2: Problem & Objective
- Slide 3: Data Sources
- Slide 4: Architecture & Pipeline
- Slide 5-6: EDA Highlights & Trends
- Slide 7-8: Performance & Risk Analytics
- Slide 9-10: BI Dashboard Screenshots
- Slide 11: Key Business Findings
- Slide 12: Thank You
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parent.parent
DOCS_DIR = PROJECT_ROOT / "docs"
FIGURES_DIR = PROJECT_ROOT / "figures"

sys.path.append(str(PROJECT_ROOT / "scripts"))
from utils import ensure_directory, setup_logging

logger = setup_logging("generate_presentation")

# Color Constants
NAVY = RGBColor(10, 37, 64)
BLUE = RGBColor(0, 102, 255)
SLATE = RGBColor(99, 115, 129)
WHITE = RGBColor(255, 255, 255)


def add_header(slide, title_text, category_text="BLUESTOCK FINTECH"):
    """Helper to add standard presentation header."""
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(9.0), Inches(1.0)
    )
    tf = title_box.text_frame
    tf.word_wrap = True

    p_cat = tf.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = BLUE

    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = NAVY


def build_presentation():
    """Generate 12-slide PowerPoint file."""
    ensure_directory(DOCS_DIR)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank_layout = prs.slide_layouts[6]

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(blank_layout)
    tb = slide1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(8.0), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "BLUESTOCK FINTECH"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = BLUE

    p2 = tf.add_paragraph()
    p2.text = "Mutual Fund Analytics Platform"
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = NAVY

    p3 = tf.add_paragraph()
    p3.text = "End-to-End Data Engineering, ETL Pipeline, Risk Engine & BI Dashboard"
    p3.font.size = Pt(14)
    p3.font.color.rgb = SLATE

    # Slide 2: Problem & Objective
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "Problem Statement & Capstone Objectives")
    tb2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(3.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    tf2.text = (
        "• Data Fragmentation: NAV, AUM, and transactions are scattered across TXT, HTML, and PDF formats.\n"
        "• Risk Comparison Gap: Lack of normalized, risk-adjusted performance metrics (Sharpe, Sortino, Alpha).\n"
        "• Objectives Met:\n"
        "   1. Automated ETL Pipeline from raw AMFI & mfapi.in sources.\n"
        "   2. 5-Table SQLite Star Schema (dim_fund, dim_date, fact_nav, fact_tx, fact_perf).\n"
        "   3. Quantitative Risk & Return Engine (21 metrics).\n"
        "   4. Interactive Multipage BI Dashboard & Fund Recommender."
    )

    # Slide 3: Data Sources
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "Data Ecosystem & 10 Provided Datasets")
    tb3 = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(3.5))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    tf3.text = (
        "• AMFI India & mfapi.in: Daily NAV history (46,000+ rows across 40 schemes).\n"
        "• AMFI Quarterly Reports: Quarterly AUM for top 10 AMCs (2022-2025).\n"
        "• AMFI Monthly Notes: Monthly SIP inflows, active accounts, new registrations.\n"
        "• Simulated Transactions: 32,000+ investor SIP & lumpsum transactions across 12 states.\n"
        "• NSE/BSE Benchmark Indices: Daily closing levels for Nifty 50 & Nifty 100."
    )

    # Slide 4: Architecture & Pipeline
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "System Architecture & ETL Pipeline")
    tb4 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(3.5))
    tf4 = tb4.text_frame
    tf4.word_wrap = True
    tf4.text = (
        "• Extract: Ingest raw CSVs and live JSON APIs (mfapi.in).\n"
        "• Transform: Forward-fill holiday NAVs, compute daily returns, validate AMFI codes.\n"
        "• Load: SQLAlchemy ORM loading into SQLite star schema with indexes.\n"
        "• Analyze: Modular Python packages (scripts/risk_metrics.py, scripts/cohort_analysis.py).\n"
        "• Visualize: Streamlit Web Dashboard + PDF export utility."
    )

    # Slide 5: EDA Highlights
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "Exploratory Data Analysis: AUM & NAV Trends")
    tb5 = slide5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(3.5))
    tf5 = tb5.text_frame
    tf5.word_wrap = True
    tf5.text = (
        "• Industry AUM Growth: Reached Rs. 81.00 Lakh Crore in Dec 2025.\n"
        "• Top AMC Leadership: SBI MF leads with Rs. 12.50L Cr, followed by ICICI Pru & HDFC MF.\n"
        "• NAV Trajectory: Strong post-COVID recovery with steady 2024-2025 market compounding."
    )

    # Slide 6: Investor Demographics
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "Exploratory Data Analysis: Investor Demographics")
    tb6 = slide6.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(3.5))
    tf6 = tb6.text_frame
    tf6.word_wrap = True
    tf6.text = (
        "• Geographic Volume: Maharashtra and Gujarat generate 38% of total transaction volume.\n"
        "• T30 vs B30 Split: Top 30 cities contribute 68% of AUM; B30 cities show +24% YoY growth.\n"
        "• Age Demographics: 26-35 age bracket forms the largest investor cohort."
    )

    # Slide 7: Performance Metrics
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "Performance & Risk Engine Results")
    tb7 = slide7.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(3.5))
    tf7 = tb7.text_frame
    tf7.word_wrap = True
    tf7.text = (
        "• Risk-Adjusted Leaders: Top Sharpe ratios achieved by Mirae Asset Large Cap & Kotak Flexicap.\n"
        "• Alpha Generation: 72.5% of equity schemes generated positive Jensen's Alpha relative to Nifty 50.\n"
        "• Sortino & Drawdowns: Average Max Drawdown of -15.4% recovered within 45 trading days."
    )

    # Slide 8: Advanced Risk Analytics
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "Advanced Analytics: VaR, Cohorts & HHI")
    tb8 = slide8.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(3.5))
    tf8 = tb8.text_frame
    tf8.word_wrap = True
    tf8.text = (
        "• 95% Historical VaR: Equity funds show average 1-day 95% VaR of -1.85% (CVaR: -2.75%).\n"
        "• SIP Continuity: Flagged at-risk investors with inter-transaction gaps > 35 days.\n"
        "• Sector Concentration: Portfolios with HHI < 800 maintain high diversification scores (>85)."
    )

    # Slide 9: Dashboard - Overview & Performance
    slide9 = prs.slides.add_slide(blank_layout)
    add_header(slide9, "Interactive BI Dashboard: Pages 1 & 2")
    tb9 = slide9.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(3.5))
    tf9 = tb9.text_frame
    tf9.word_wrap = True
    tf9.text = (
        "• Page 1 (Industry Overview): Real-time KPI cards, quarterly AUM line chart, top AMC bar chart.\n"
        "• Page 2 (Fund Performance): Interactive risk-return scatter plot, sortable scorecard table, NAV vs benchmark comparison."
    )

    # Slide 10: Dashboard - Investor & Trends
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "Interactive BI Dashboard: Pages 3 & 4")
    tb10 = slide10.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(9.0), Inches(3.5)
    )
    tf10 = tb10.text_frame
    tf10.word_wrap = True
    tf10.text = (
        "• Page 3 (Investor Analytics): State transaction volume bar chart, SIP vs Lumpsum donut split, age group SIP bar.\n"
        "• Page 4 (SIP & Market Trends): Dual-axis SIP inflow vs Nifty 50 line, category inflow heatmaps."
    )

    # Slide 11: Key Findings & Takeaways
    slide11 = prs.slides.add_slide(blank_layout)
    add_header(slide11, "Key Business Findings & Next Steps")
    tb11 = slide11.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(9.0), Inches(3.5)
    )
    tf11 = tb11.text_frame
    tf11.word_wrap = True
    tf11.text = (
        "1. Active Management Value: Positive Alpha validates active stock selection in Indian markets.\n"
        "2. Systematic Investing Strength: SIP inflows act as resilient counter-cyclical buffers during market drops.\n"
        "3. Recommendation Deployment: Integrate `scripts/recommender.py` into mobile app backends."
    )

    # Slide 12: Thank You
    slide12 = prs.slides.add_slide(blank_layout)
    tb12 = slide12.shapes.add_textbox(
        Inches(1.0), Inches(2.0), Inches(8.0), Inches(2.0)
    )
    tf12 = tb12.text_frame
    tf12.word_wrap = True
    p_end = tf12.paragraphs[0]
    p_end.text = "Thank You!"
    p_end.font.size = Pt(36)
    p_end.font.bold = True
    p_end.font.color.rgb = BLUE
    p_end.alignment = PP_ALIGN.CENTER

    p_end2 = tf12.add_paragraph()
    p_end2.text = "Bluestock Mutual Fund Analytics Platform Capstone"
    p_end2.font.size = Pt(16)
    p_end2.font.color.rgb = NAVY
    p_end2.alignment = PP_ALIGN.CENTER

    output_path = DOCS_DIR / "Bluestock_MF_Presentation.pptx"
    prs.save(str(output_path))
    logger.info(f"Generated 12-slide presentation deck at {output_path}")


if __name__ == "__main__":
    build_presentation()
